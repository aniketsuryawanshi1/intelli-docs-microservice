import pandas as pd
from pathlib import Path
from pptx import Presentation
from PIL import Image
import pytesseract
from docx import Document as DocxDocument
from unstructured.partition.auto import partition

def split_text(text, max_length=300):
    sentences = text.split('. ')
    chunks, current = [], ""
    for sent in sentences:
        if len(current) + len(sent) < max_length:
            current += sent + ". "
        else:
            chunks.append(current.strip())
            current = sent + ". "
    if current:
        chunks.append(current.strip())
    return chunks

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text_from_image(file_path):
    image = Image.open(file_path)
    return pytesseract.image_to_string(image)

def extract_text_from_docx(file_path):
    doc = DocxDocument(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_xlsx(file_path):
    try:
        dfs = pd.read_excel(file_path, sheet_name=None)
        all_text = []
        for sheet_name, df in dfs.items():
            all_text.append(f"Sheet: {sheet_name}")
            for _, row in df.iterrows():
                row_data = [f"{col}: {val}" for col, val in row.items()]
                sentence = ", ".join(row_data)
                all_text.append(sentence)
        return "\n".join(all_text)
    except Exception as e:
        print(f"Error reading Excel: {e}")
        return ""

def extract_text_from_csv(file_path):
    try:
        df = pd.read_csv(file_path)
        cols = list(df.columns)
        all_text = [f"Table columns: {', '.join(cols)}."]
        for _, row in df.iterrows():
            row_data = [f"{col}: {val}" for col, val in row.items()]
            sentence = ", ".join(row_data)
            all_text.append(sentence)
        return "\n".join(all_text)
    except Exception as e:
        print(f"Error reading CSV: {e}")
        return ""

def process_file(file_path):
    suffix = Path(file_path).suffix.lower()
    text = ""
    try:
        # Use unstructured.io for advanced parsing
        elements = partition(file_path)
        text = "\n".join([el.text for el in elements if el.text is not None])
        if not text.strip():
            raise ValueError("Empty text extracted.")
    except Exception as e:
        print(f"Unstructured failed for {file_path}: {e}")

    if not text.strip():
        try:
            if suffix == ".docx":
                text = extract_text_from_docx(file_path)
            elif suffix == ".xlsx":
                text = extract_text_from_xlsx(file_path)
            elif suffix == ".csv":
                text = extract_text_from_csv(file_path)
            elif suffix in [".png", ".jpg", ".jpeg", ".tiff"]:
                text = extract_text_from_image(file_path)
            elif suffix == ".pptx":
                text = extract_text_from_pptx(file_path)
        except Exception as e:
            print(f"Fallback error for {suffix}: {e}")
            raise ValueError(f"Could not process the file: {file_path}")

    if not text.strip():
        raise ValueError("No extractable text found.")
    return split_text(text)